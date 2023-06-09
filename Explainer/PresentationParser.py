from pptx import Presentation


class PresentationParser:
    def __init__(self, presentation_path):
        self.presentation_path = presentation_path

    def process_presentation(self) -> list[str]:
        """
        Process each slide in the presentation asynchronously and retrieve the responses.
        :return: list[str]: List of explanations for each slide.
        """
        prs = Presentation(self.presentation_path)
        slides = []
        for slide in prs.slides:
            slides.append(self.process_slide_text(slide))
        return slides

    @staticmethod
    def process_slide_text(slide) -> str:
        """
        Extract the text content from a slide in the PowerPoint presentation.
        :param slide: Slide object from the PowerPoint presentation.
        :return: str: Processed text content from the slide.
        """
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text.append(run.text.strip())
        return " ".join(slide_text)
